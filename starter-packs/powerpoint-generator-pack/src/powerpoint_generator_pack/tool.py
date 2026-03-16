"""Generate PowerPoint presentations with proper layouts."""

from __future__ import annotations

import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# Standard slide layout indices in the default template
_LAYOUT_TITLE_SLIDE = 0       # Title Slide
_LAYOUT_TITLE_CONTENT = 1     # Title and Content
_LAYOUT_SECTION_HEADER = 2    # Section Header
_LAYOUT_TWO_CONTENT = 3       # Two Content
_LAYOUT_BLANK = 6             # Blank


def _add_title_slide(prs: Presentation, title: str, content: str) -> None:
    """Add a title slide with title and subtitle."""
    layout = prs.slide_layouts[_LAYOUT_TITLE_SLIDE]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = content
    else:
        # Fallback: add a text box for the subtitle
        from pptx.util import Emu
        txBox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.5))
        txBox.text_frame.text = content


def _add_content_slide(prs: Presentation, title: str, content: str) -> None:
    """Add a content slide with title and bullet points."""
    layout = prs.slide_layouts[_LAYOUT_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    # Content placeholder is usually index 1
    body_placeholder = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_placeholder = ph
            break

    if body_placeholder is not None:
        tf = body_placeholder.text_frame
        tf.clear()
        lines = content.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                tf.text = line.strip()
            else:
                p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 0
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = content


def _add_two_column_slide(prs: Presentation, title: str, content: str) -> None:
    """Add a two-column slide. Content is split by '|||' or newline-heavy blocks."""
    layout_idx = _LAYOUT_TWO_CONTENT
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError:
        # Fallback to content layout
        _add_content_slide(prs, title, content)
        return

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    # Split content for two columns
    if "|||" in content:
        left_text, right_text = content.split("|||", 1)
    else:
        lines = content.split("\n")
        mid = len(lines) // 2
        left_text = "\n".join(lines[:mid])
        right_text = "\n".join(lines[mid:])

    # Fill placeholders (indices 1 and 2 for two-content layout)
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    if 1 in placeholders:
        placeholders[1].text_frame.text = left_text.strip()
    if 2 in placeholders:
        placeholders[2].text_frame.text = right_text.strip()


def _add_blank_slide(prs: Presentation, title: str, content: str) -> None:
    """Add a blank slide with optional text box."""
    layout_idx = _LAYOUT_BLANK
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError:
        layout = prs.slide_layouts[len(prs.slide_layouts) - 1]

    slide = prs.slides.add_slide(layout)

    # Add a text box if there is content
    combined = f"{title}\n\n{content}".strip() if title or content else ""
    if combined:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = combined


_LAYOUT_HANDLERS = {
    "title": _add_title_slide,
    "content": _add_content_slide,
    "two_column": _add_two_column_slide,
    "blank": _add_blank_slide,
}


def run(
    slides: list[dict],
    output_path: str = "",
    title: str = "Presentation",
) -> dict:
    """Create a PowerPoint presentation from a list of slide definitions.

    Args:
        slides: List of dicts, each with keys 'title', 'content', and 'layout'
                (one of 'title', 'content', 'two_column', 'blank').
        output_path: Where to save the .pptx file. Auto-generated if empty.
        title: Presentation title used in file naming when output_path is empty.

    Returns:
        Dictionary with output_path and slide_count.
    """
    if not slides:
        raise ValueError("At least one slide is required.")

    prs = Presentation()

    for slide_def in slides:
        slide_title = slide_def.get("title", "")
        slide_content = slide_def.get("content", "")
        layout = slide_def.get("layout", "content")

        handler = _LAYOUT_HANDLERS.get(layout)
        if handler is None:
            raise ValueError(
                f"Unknown layout '{layout}'. "
                f"Choose from: {', '.join(_LAYOUT_HANDLERS)}"
            )
        handler(prs, slide_title, slide_content)

    # Determine output path
    if not output_path:
        safe_title = "".join(
            c if c.isalnum() or c in (" ", "-", "_") else "_" for c in title
        ).strip()
        output_path = os.path.join(
            tempfile.gettempdir(), f"{safe_title or 'presentation'}.pptx"
        )

    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    return {
        "output_path": output_path,
        "slide_count": len(slides),
    }
